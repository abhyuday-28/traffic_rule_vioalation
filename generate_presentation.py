from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
OUTPUT = ROOT / "Traffic_Violation_System_Presentation.pptx"


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(15, 23, 42)

    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.3))
    band.fill.solid()
    band.fill.fore_color.rgb = RGBColor(30, 64, 175)
    band.line.color.rgb = RGBColor(30, 64, 175)

    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.7), Inches(11.5), Inches(1.4))
    p = title_box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = RGBColor(255, 255, 255)

    subtitle_box = slide.shapes.add_textbox(Inches(0.7), Inches(3.0), Inches(10.5), Inches(1.5))
    p = subtitle_box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = subtitle
    r.font.size = Pt(17)
    r.font.color.rgb = RGBColor(226, 232, 240)


def add_bullets_slide(prs: Presentation, title: str, bullets: list[str], note: str | None = None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(248, 250, 252)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(11.5), Inches(0.7))
    p = title_box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = RGBColor(15, 23, 42)

    body = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.0), Inches(5.6))
    tf = body.text_frame
    tf.word_wrap = True
    for index, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(30, 41, 59)
        p.space_after = Pt(12)

    if note:
        note_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.7), Inches(10.8), Inches(0.5))
        p = note_box.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = note
        r.font.size = Pt(12)
        r.font.italic = True
        r.font.color.rgb = RGBColor(100, 116, 139)


def add_two_column_slide(prs: Presentation, title: str, left_title: str, left_items: list[str], right_title: str, right_items: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(11.5), Inches(0.7))
    p = title_box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = RGBColor(15, 23, 42)

    for x, heading, items, fill in (
        (0.6, left_title, left_items, RGBColor(239, 246, 255)),
        (6.2, right_title, right_items, RGBColor(248, 250, 252)),
    ):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.2), Inches(5.0), Inches(5.5))
        card.fill.solid()
        card.fill.fore_color.rgb = fill
        card.line.color.rgb = RGBColor(203, 213, 225)

        box = slide.shapes.add_textbox(Inches(x + 0.25), Inches(1.45), Inches(4.5), Inches(5.0))
        tf = box.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = heading
        r.font.size = Pt(19)
        r.font.bold = True
        r.font.color.rgb = RGBColor(30, 64, 175)
        for item in items:
            p = tf.add_paragraph()
            p.text = item
            p.level = 0
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(30, 41, 59)
            p.space_after = Pt(8)


def add_flow_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(248, 250, 252)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(11.5), Inches(0.7))
    p = title_box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "System Workflow"
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = RGBColor(15, 23, 42)

    steps = [
        "Input source\nImage / Video / Camera",
        "Traffic detection\nYOLOX + helmet detector",
        "Rule engine\nNo helmet / Triple riding / Red light",
        "Plate stage\nWPOD-NET + OCR fallback",
        "Evidence & action\nLog, preview, email authority",
    ]
    left = 0.45
    for idx, step in enumerate(steps):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left + idx * 2.28), Inches(2.3), Inches(2.0), Inches(1.4))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(219, 234, 254) if idx % 2 == 0 else RGBColor(224, 231, 255)
        box.line.color.rgb = RGBColor(59, 130, 246)
        tf = box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = step
        r.font.size = Pt(15)
        r.font.bold = True
        r.font.color.rgb = RGBColor(30, 41, 59)
        if idx < len(steps) - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(left + idx * 2.28 + 1.95), Inches(2.72), Inches(0.35), Inches(0.55))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = RGBColor(37, 99, 235)
            arrow.line.color.rgb = RGBColor(37, 99, 235)


def build_presentation() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "AI-Based Traffic Rule Violation Detection System",
        "Project presentation covering problem statement, system design, current implementation, results, and future scope.",
    )

    add_bullets_slide(
        prs,
        "Problem Statement",
        [
            "Manual traffic monitoring is slow, labor-intensive, and inconsistent.",
            "Police cannot reliably track multiple violations at the same time from live traffic feeds.",
            "Violations such as no helmet, triple riding, and red-light jumping need evidence and vehicle identification.",
            "A software system can reduce manual effort and create faster violation reporting workflows.",
        ],
    )

    add_bullets_slide(
        prs,
        "Project Objective",
        [
            "Build a desktop application that accepts image, video, or live camera input.",
            "Detect major two-wheeler traffic violations automatically.",
            "Extract the vehicle number plate and store evidence images.",
            "Allow the operator to review the violation and send an email report to traffic authorities.",
        ],
    )

    add_flow_slide(prs)

    add_two_column_slide(
        prs,
        "Core Modules",
        "Detection Modules",
        [
            "Traffic object detection using YOLOX through OpenCV DNN.",
            "Helmet detection using custom YOLO weights.",
            "Rule-based rider-to-bike grouping for no-helmet and triple-riding analysis.",
            "Traffic light state estimation for red-light jumping.",
        ],
        "Support Modules",
        [
            "Plate extraction with WPOD-NET and OCR fallback.",
            "Tkinter GUI for operator workflow.",
            "Evidence image saving and CSV log export.",
            "SMTP email integration for authority notification.",
        ],
    )

    add_bullets_slide(
        prs,
        "User Interface",
        [
            "Single-window Tkinter application for easy operator use.",
            "Supports image upload, video file input, and live camera capture.",
            "Displays annotated frames, detected violations, and number plate preview.",
            "Provides export and email actions directly from the dashboard.",
        ],
    )

    add_bullets_slide(
        prs,
        "Implemented Violations",
        [
            "No Helmet: bike-level rule comparing assigned riders and matched helmets.",
            "Triple Riding: flags a bike when rider count on the same bike is three or more.",
            "Red-Light Jumping: vehicle crossing the configured stop line while the signal is red.",
            "Number Plate Extraction: attempts plate crop and OCR for each detected violation.",
        ],
        note="Current implementation is a working prototype and still depends on scene quality, camera angle, and model behavior.",
    )

    add_two_column_slide(
        prs,
        "Technology Stack",
        "Software",
        [
            "Python",
            "OpenCV",
            "TensorFlow / Keras",
            "Tkinter",
            "Pandas",
            "SMTP email",
        ],
        "Models",
        [
            "YOLOX ONNX for traffic objects",
            "Custom helmet detector",
            "WPOD-NET for plate extraction",
            "Character recognition OCR model",
        ],
    )

    add_bullets_slide(
        prs,
        "Current Results",
        [
            "The application can process images, videos, and live camera input from one GUI.",
            "It creates evidence images, extracts a number plate image when possible, and logs violations.",
            "Helmet, triple-riding, and red-light logic are integrated into one workflow.",
            "Email reporting is available from the operator interface after a violation is selected.",
        ],
    )

    add_bullets_slide(
        prs,
        "Challenges Faced",
        [
            "Crowded bikes and overlapping riders are hard to count reliably.",
            "Helmet detection can confuse hair, shadows, or low-quality regions.",
            "OCR accuracy depends heavily on plate visibility, blur, and crop quality.",
            "Raspberry Pi deployment is difficult without simplifying the pipeline.",
        ],
    )

    add_bullets_slide(
        prs,
        "Future Enhancements",
        [
            "Train a dedicated violation model instead of combining multiple heuristics.",
            "Improve OCR with stronger plate-specific recognition.",
            "Add a debug calibration mode for rider, head, and helmet counts.",
            "Deploy a lighter version for edge devices or use hardware acceleration.",
            "Connect to authority databases for automatic challan generation.",
        ],
    )

    add_bullets_slide(
        prs,
        "Conclusion",
        [
            "The project demonstrates an end-to-end AI-assisted traffic violation workflow.",
            "It combines detection, rule evaluation, plate extraction, evidence capture, and authority reporting.",
            "The current system is suitable as a prototype and academic demonstration.",
            "With better training data and stronger models, it can be pushed toward practical deployment.",
        ],
    )

    prs.save(OUTPUT)
    return OUTPUT


if __name__ == "__main__":
    path = build_presentation()
    print(path)
